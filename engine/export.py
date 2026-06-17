#!/usr/bin/env python3
"""PitchGraph export — turn a finished .pptx into deliverables. Engine-agnostic
(works for every build path).

  pdf      <deck.pptx> [out.pdf]                 -> PDF (LibreOffice)
  handout  <deck.pptx> [out.pdf] [--cols 2 --rows 2]  -> N-up handout PDF
  notes    <deck.pptx> <brief.json>              -> inject per-slide speaker notes
  all      <deck.pptx> <brief.json>              -> notes + pdf + 2x2 handout

Speaker notes come from each brief outline item's `notes` field (the researcher /
note-writer fills them). Injection uses python-pptx so it applies to decks from
ANY engine, not just the from-scratch one.
"""
from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
from pathlib import Path


def _soffice():
    return shutil.which("soffice") or (
        "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if Path("/Applications/LibreOffice.app/Contents/MacOS/soffice").exists() else None)


def to_pdf(deck: Path, out: Path | None = None) -> Path:
    so = _soffice()
    if not so:
        raise RuntimeError("LibreOffice (soffice) not found — needed for PDF export.")
    outdir = (out.parent if out else deck.parent)
    outdir.mkdir(parents=True, exist_ok=True)
    subprocess.run([so, "--headless", "--convert-to", "pdf", "--outdir", str(outdir), str(deck)],
                   check=True, capture_output=True)
    made = outdir / (deck.stem + ".pdf")
    if out and out != made:
        shutil.move(str(made), str(out))
        return out
    return made


def inject_notes(deck: Path, brief: dict) -> int:
    from pptx import Presentation
    prs = Presentation(str(deck))
    outline = brief.get("outline", [])
    n = 0
    for slide, item in zip(prs.slides, outline):
        note = item.get("notes")
        if note:
            slide.notes_slide.notes_text_frame.text = note
            n += 1
    prs.save(str(deck))
    return n


def handout(deck: Path, out: Path | None, cols: int = 2, rows: int = 3) -> Path:
    """N-up handout via Chromium print-to-PDF (no Pillow PDF/JPEG dependency)."""
    import base64
    pdf = to_pdf(deck)
    tmp = deck.parent / (deck.stem + "_handout_pages")
    tmp.mkdir(exist_ok=True)
    subprocess.run(["pdftoppm", "-png", "-r", "110", str(pdf), str(tmp / "p")],
                   check=True, capture_output=True)
    pages = sorted(tmp.glob("p*.png"))
    if not pages:
        raise RuntimeError("no slide images rendered for handout")
    imgs = []
    for p in pages:
        b64 = base64.b64encode(p.read_bytes()).decode()
        imgs.append(f'<img src="data:image/png;base64,{b64}">')
    html = f"""<!doctype html><html><head><style>
      @page {{ size: A4; margin: 12mm; }}
      body {{ margin:0; }}
      .grid {{ display:grid; grid-template-columns: repeat({cols}, 1fr); gap: 8mm; }}
      img {{ width:100%; border:1px solid #ddd; }}
    </style></head><body><div class="grid">{''.join(imgs)}</div></body></html>"""
    out = out or deck.with_name(deck.stem + "_handout.pdf")
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as pw:
            browser = pw.chromium.launch()
            page = browser.new_page()
            page.set_content(html, wait_until="load")
            page.pdf(path=str(out), format="A4", print_background=True)
            browser.close()
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    return out


def main(argv):
    ap = argparse.ArgumentParser(description="PitchGraph export")
    sub = ap.add_subparsers(dest="cmd", required=True)
    p1 = sub.add_parser("pdf"); p1.add_argument("deck"); p1.add_argument("out", nargs="?")
    p2 = sub.add_parser("handout"); p2.add_argument("deck"); p2.add_argument("out", nargs="?")
    p2.add_argument("--cols", type=int, default=2); p2.add_argument("--rows", type=int, default=2)
    p3 = sub.add_parser("notes"); p3.add_argument("deck"); p3.add_argument("brief")
    p4 = sub.add_parser("all"); p4.add_argument("deck"); p4.add_argument("brief")
    a = ap.parse_args(argv[1:])
    deck = Path(a.deck)

    if a.cmd == "pdf":
        print("PDF:", to_pdf(deck, Path(a.out) if a.out else None))
    elif a.cmd == "handout":
        print("Handout:", handout(deck, Path(a.out) if a.out else None, a.cols, a.rows))
    elif a.cmd == "notes":
        n = inject_notes(deck, json.loads(Path(a.brief).read_text()))
        print(f"Injected speaker notes into {n} slides of {deck}")
    elif a.cmd == "all":
        brief = json.loads(Path(a.brief).read_text())
        n = inject_notes(deck, brief)
        pdf = to_pdf(deck)
        ho = handout(deck, None)
        print(f"notes:{n} slides | pdf:{pdf} | handout:{ho}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
