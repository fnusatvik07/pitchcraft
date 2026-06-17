"""The nine native layout builders for PitchCraft.

Every builder has the signature `(slide, content, theme, ctx) -> None` and reads
ALL styling from `theme` -- never a hardcoded color/font/size. `ctx` carries
shared helpers (asset placement, chart builder) so layouts stay engine-agnostic.

Slide canvas is 16:9 = 13.333in x 7.5in.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from theme import Theme, hex_to_rgb

EMU_PER_IN = 914400
SLIDE_W = 13.333
SLIDE_H = 7.5

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


# --------------------------------------------------------------------------- #
# low-level native-shape helpers
# --------------------------------------------------------------------------- #
def fill_background(slide, theme: Theme, key: str = "bg"):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.color(key)


def add_rect(slide, x, y, w, h, color, rounded=False, radius_pt=6, line_color=None, line_w=0.75):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:  # adjustment[0] is the corner radius as a fraction of the shorter side
            shorter_pt = min(w, h) * 72.0
            shp.adjustments[0] = max(0.0, min(0.5, radius_pt / shorter_pt))
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    runs,
    *,
    align="left",
    anchor=MSO_ANCHOR.TOP,
    wrap=True,
):
    """runs: list of paragraphs; each paragraph is a list of (text, opts) tuples.

    opts keys: font, size (Pt), color (RGBColor), bold, italic, space_after (Pt),
    line_spacing (float), bullet (str prefix).
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN[align]
        first_opts = para[0][1] if para else {}
        if "space_after" in first_opts:
            p.space_after = first_opts["space_after"]
        if "line_spacing" in first_opts:
            p.line_spacing = first_opts["line_spacing"]
        prefix = first_opts.get("bullet")
        for j, (text, opts) in enumerate(para):
            r = p.add_run()
            r.text = (prefix + text) if (prefix and j == 0) else text
            f = r.font
            f.name = opts.get("font")
            if opts.get("size"):
                f.size = opts["size"]
            if opts.get("color") is not None:
                f.color.rgb = opts["color"]
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
    return tb


def add_caption(slide, x, y, w, theme: Theme, text):
    """A small muted source/credit caption."""
    add_text(
        slide, x, y, w, 0.3,
        [[(text, {"font": theme.font("body"), "size": theme.size("caption"),
                  "color": theme.color("muted"), "italic": True})]],
    )


def _img_region(slide, content, theme, ctx, x, y, w, h):
    """Place an image if present; otherwise a token-colored placeholder block."""
    img = content.get("image")
    if img and ctx and ctx.get("place_image"):
        ok = ctx["place_image"](slide, img, x, y, w, h, theme)
        if ok:
            return
    # graceful fallback: a surface block with a subtle label
    add_rect(slide, x, y, w, h, theme.color("surface"),
             rounded=True, radius_pt=theme.layout.get("corner_radius_pt", 6))
    add_text(slide, x, y + h / 2 - 0.2, w, 0.4,
             [[("image", {"font": theme.font("body"), "size": theme.size("caption"),
                          "color": theme.color("muted")})]], align="center")


# --------------------------------------------------------------------------- #
# the nine layouts
# --------------------------------------------------------------------------- #
def title(slide, content, theme, ctx=None):
    fill_background(slide, theme, "bg")
    align = theme.layout.get("title_align", "left")
    x = theme.margin
    w = SLIDE_W - 2 * theme.margin
    if theme.layout.get("section_style") == "band":
        add_rect(slide, 0, SLIDE_H - 0.18, SLIDE_W, 0.18, theme.color("primary"))
    runs = []
    if content.get("kicker"):
        runs.append([(content["kicker"].upper(), {
            "font": theme.font("body"), "size": theme.size("caption"),
            "color": theme.color("accent"), "bold": True, "space_after": Pt(10)})])
    runs.append([(content["title"], {
        "font": theme.font("display"), "size": theme.size("title"),
        "color": theme.color("ink"), "bold": True, "line_spacing": 1.04,
        "space_after": Pt(10)})])
    if content.get("subtitle"):
        runs.append([(content["subtitle"], {
            "font": theme.font("body"), "size": theme.size("h2"),
            "color": theme.color("muted"), "line_spacing": 1.1})])
    add_text(slide, x, 2.4, w, 3.0, runs, align=align, anchor=MSO_ANCHOR.TOP)


def agenda(slide, content, theme, ctx=None):
    fill_background(slide, theme, "bg")
    _heading(slide, theme, content.get("title", "Agenda"))
    items = content.get("items", [])
    top = 2.0
    row_h = min(0.95, (SLIDE_H - top - theme.margin) / max(len(items), 1))
    bullet = theme.layout.get("bullet_char", "—")
    for i, it in enumerate(items):
        y = top + i * row_h
        add_text(slide, theme.margin, y, 0.8, row_h,
                 [[(f"{i + 1:02d}", {"font": theme.font("display"),
                                     "size": theme.size("h2"), "color": theme.color("accent"),
                                     "bold": True})]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, theme.margin + 0.9, y, SLIDE_W - 2 * theme.margin - 0.9, row_h,
                 [[(it, {"font": theme.font("body"), "size": theme.size("h2"),
                         "color": theme.color("ink")})]], anchor=MSO_ANCHOR.MIDDLE)


def section(slide, content, theme, ctx=None):
    style = theme.layout.get("section_style", "band")
    if style == "fullbleed":
        fill_background(slide, theme, "primary")
        ink, accent = theme.color("bg"), theme.color("accent")
    else:
        fill_background(slide, theme, "bg")
        add_rect(slide, 0, 2.9, SLIDE_W, 1.7, theme.color("surface"))
        ink, accent = theme.color("ink"), theme.color("accent")
    num = content.get("number")
    runs = []
    if num:
        runs.append([(f"{num:02d}" if isinstance(num, int) else str(num), {
            "font": theme.font("display"), "size": Pt(theme.type_scale["title"] + 8),
            "color": accent, "bold": True, "space_after": Pt(6)})])
    runs.append([(content["title"], {
        "font": theme.font("display"), "size": theme.size("title"),
        "color": ink, "bold": True})])
    add_text(slide, theme.margin, 3.0, SLIDE_W - 2 * theme.margin, 1.6, runs,
             align=theme.layout.get("title_align", "left"), anchor=MSO_ANCHOR.MIDDLE)


def content_image(slide, content, theme, ctx=None):
    fill_background(slide, theme, "bg")
    _heading(slide, theme, content.get("title", ""))
    top = 2.0
    h = SLIDE_H - top - theme.margin
    img = content.get("image")
    img_path = (img.get("path") if isinstance(img, dict) else img) if img else None
    if not img_path:
        # no image -> bullets span the full width (no empty placeholder)
        _bullets(slide, theme, content.get("bullets", []), theme.margin, top,
                 SLIDE_W - 2 * theme.margin, h)
        return
    img_left = content.get("image_side", "right") == "left"
    col_w = (SLIDE_W - 2 * theme.margin - theme.gutter) / 2
    text_x = theme.margin + (0 if not img_left else col_w + theme.gutter)
    img_x = theme.margin + (col_w + theme.gutter if not img_left else 0)
    _bullets(slide, theme, content.get("bullets", []), text_x, top, col_w, h)
    _img_region(slide, content, theme, ctx, img_x, top, col_w, h)


def stat_callout(slide, content, theme, ctx=None):
    fill_background(slide, theme, "bg")
    align = theme.layout.get("title_align", "left")
    add_text(slide, theme.margin, 2.3, SLIDE_W - 2 * theme.margin, 2.0,
             [[(content["stat"], {"font": theme.font("display"),
                                  "size": theme.size("stat"), "color": theme.color("primary"),
                                  "bold": True})]], align=align)
    add_text(slide, theme.margin, 4.5, SLIDE_W - 2 * theme.margin, 1.2,
             [[(content.get("label", ""), {"font": theme.font("body"),
                                           "size": theme.size("h2"), "color": theme.color("ink")})]],
             align=align)
    src = content.get("source")
    if src:
        add_caption(slide, theme.margin, SLIDE_H - theme.margin - 0.1,
                    SLIDE_W - 2 * theme.margin, theme, _src_text(src))


def comparison(slide, content, theme, ctx=None):
    fill_background(slide, theme, "bg")
    _heading(slide, theme, content.get("title", ""))
    cols = content.get("columns", [])[:2]
    col_w = (SLIDE_W - 2 * theme.margin - theme.gutter) / 2
    top, h = 2.0, SLIDE_H - 2.0 - theme.margin
    accents = ["primary", "accent"]
    for i, col in enumerate(cols):
        x = theme.margin + i * (col_w + theme.gutter)
        card = add_rect(slide, x, top, col_w, h, theme.color("surface"),
                        rounded=True, radius_pt=theme.layout.get("corner_radius_pt", 6))
        add_rect(slide, x, top, col_w, 0.12, theme.color(accents[i % 2]))
        add_text(slide, x + 0.3, top + 0.3, col_w - 0.6, 0.6,
                 [[(col.get("heading", ""), {"font": theme.font("display"),
                                             "size": theme.size("h2"),
                                             "color": theme.color("ink"), "bold": True})]])
        _bullets(slide, theme, col.get("points", []), x + 0.3, top + 1.1, col_w - 0.6, h - 1.4)


def quote(slide, content, theme, ctx=None):
    fill_background(slide, theme, "bg")
    add_rect(slide, theme.margin, 2.2, 0.08, 2.6, theme.color("accent"))
    add_text(slide, theme.margin + 0.4, 2.2, SLIDE_W - 2 * theme.margin - 0.4, 2.6,
             [[("“" + content["quote"] + "”", {
                 "font": theme.font("display"), "size": theme.size("h2"),
                 "color": theme.color("ink"), "italic": True, "line_spacing": 1.15})]],
             anchor=MSO_ANCHOR.MIDDLE)
    if content.get("attribution"):
        add_text(slide, theme.margin + 0.4, 5.0, SLIDE_W - 2 * theme.margin - 0.4, 0.5,
                 [[("— " + content["attribution"], {
                     "font": theme.font("body"), "size": theme.size("body"),
                     "color": theme.color("muted"), "bold": True})]])


def chart(slide, content, theme, ctx=None):
    fill_background(slide, theme, "bg")
    _heading(slide, theme, content.get("title", ""))
    top = 2.0
    h = SLIDE_H - top - theme.margin - (0.4 if content.get("source") else 0)
    if ctx and ctx.get("add_chart"):
        ctx["add_chart"](slide, content, theme,
                         theme.margin, top, SLIDE_W - 2 * theme.margin, h)
    if content.get("source"):
        add_caption(slide, theme.margin, SLIDE_H - theme.margin - 0.1,
                    SLIDE_W - 2 * theme.margin, theme, _src_text(content["source"]))


def closing(slide, content, theme, ctx=None):
    style = theme.layout.get("section_style", "band")
    if style == "fullbleed":
        fill_background(slide, theme, "primary")
        ink = theme.color("bg")
    else:
        fill_background(slide, theme, "bg")
        ink = theme.color("ink")
    add_text(slide, theme.margin, 2.6, SLIDE_W - 2 * theme.margin, 1.4,
             [[(content.get("title", "Thank you"), {
                 "font": theme.font("display"), "size": theme.size("title"),
                 "color": ink, "bold": True})]],
             align=theme.layout.get("title_align", "left"))
    if content.get("subtitle"):
        add_text(slide, theme.margin, 4.0, SLIDE_W - 2 * theme.margin, 0.8,
                 [[(content["subtitle"], {"font": theme.font("body"),
                                          "size": theme.size("h2"),
                                          "color": ink if style == "fullbleed" else theme.color("muted")})]],
                 align=theme.layout.get("title_align", "left"))
    sources = content.get("sources", [])
    if sources:
        runs = [[(s, {"font": theme.font("body"), "size": theme.size("caption"),
                      "color": ink if style == "fullbleed" else theme.color("muted"),
                      "space_after": Pt(2)})] for s in sources[:6]]
        add_text(slide, theme.margin, 5.2, SLIDE_W - 2 * theme.margin, 2.0, runs,
                 align=theme.layout.get("title_align", "left"))


# --------------------------------------------------------------------------- #
# shared sub-helpers
# --------------------------------------------------------------------------- #
def _heading(slide, theme, text):
    add_text(slide, theme.margin, theme.margin, SLIDE_W - 2 * theme.margin, 1.1,
             [[(text, {"font": theme.font("display"), "size": theme.size("h2") if False else theme.size("h2"),
                       "color": theme.color("ink"), "bold": True})]],
             align=theme.layout.get("title_align", "left"))
    if theme.layout.get("section_style") == "band":
        add_rect(slide, theme.margin, theme.margin + 0.85, 0.6, 0.06, theme.color("accent"))


def _bullets(slide, theme, items, x, y, w, h):
    if not items:
        return
    runs = []
    bullet = theme.layout.get("bullet_char", "—") + "  "
    for it in items:
        runs.append([(it, {"font": theme.font("body"), "size": theme.size("body"),
                           "color": theme.color("ink"), "bullet": bullet,
                           "space_after": Pt(10), "line_spacing": 1.12})])
    add_text(slide, x, y, w, h, runs, anchor=MSO_ANCHOR.TOP)


def _src_text(src):
    if isinstance(src, dict):
        t = src.get("text", "")
        u = src.get("url", "")
        return "Source: " + (f"{t} — {u}" if t and u else (t or u))
    return "Source: " + str(src)


LAYOUTS = {
    "title": title,
    "agenda": agenda,
    "section": section,
    "content-image": content_image,
    "content": content_image,
    "stat-callout": stat_callout,
    "comparison": comparison,
    "quote": quote,
    "chart": chart,
    "closing": closing,
}
