"""Native (editable) chart helpers, styled by theme tokens.

Produces real Office charts via python-pptx -- not images. Series colors and
fonts come from the theme so charts never look like default Office output.
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from theme import Theme

CHART_TYPES = {
    "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "BAR": XL_CHART_TYPE.BAR_CLUSTERED,
    "LINE": XL_CHART_TYPE.LINE,
    "LINE_MARKERS": XL_CHART_TYPE.LINE_MARKERS,
    "PIE": XL_CHART_TYPE.PIE,
    "AREA": XL_CHART_TYPE.AREA,
}


def add_chart(slide, content, theme: Theme, x, y, w, h):
    ctype = CHART_TYPES.get(
        content.get("chart_type", theme.chart_style.get("type_default", "COLUMN")).upper(),
        XL_CHART_TYPE.COLUMN_CLUSTERED,
    )
    data = CategoryChartData()
    data.categories = content.get("categories", [])
    series = content.get("series", [])
    for s in series:
        data.add_series(s.get("name", ""), s.get("values", []))

    gf = slide.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart

    # Chart title: the slide already carries a descriptive heading, and an auto
    # series-name title renders dark-on-dark on dark themes. Set it explicitly only
    # when a unit label is provided, as a real colored run so the color is honored.
    unit = content.get("unit") or content.get("y_label")
    if not unit and len(series) == 1:
        unit = series[0].get("name")  # single-series: show its name as a colored title
    try:
        if unit:
            ch.has_title = True
            tf = ch.chart_title.text_frame
            tf.text = unit  # creates a real run we can color reliably
            r = tf.paragraphs[0].runs[0]
            r.font.color.rgb = theme.color("muted")
            r.font.size = Pt(theme.type_scale["caption"] + 1)
            r.font.name = theme.font("body")
            r.font.bold = False
        else:
            ch.has_title = False
    except Exception:
        pass

    # legend
    if len(series) > 1 and ctype != XL_CHART_TYPE.PIE:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        _font(ch.legend.font, theme, theme.type_scale["caption"], theme.color("muted"))
    else:
        ch.has_legend = (ctype == XL_CHART_TYPE.PIE)
        if ch.has_legend:
            ch.legend.position = XL_LEGEND_POSITION.RIGHT
            ch.legend.include_in_layout = False

    # series colors from tokens
    colors = theme.series_colors()
    try:
        if ctype == XL_CHART_TYPE.PIE:
            pts = ch.plots[0].series[0].points
            for i, pt in enumerate(pts):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        else:
            for i, plot_series in enumerate(ch.series):
                c = colors[i % len(colors)]
                if ctype in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS):
                    plot_series.format.line.color.rgb = c
                    plot_series.format.line.width = Pt(2.5)
                else:
                    plot_series.format.fill.solid()
                    plot_series.format.fill.fore_color.rgb = c
    except Exception:
        pass  # never let chart styling crash a build

    # axis fonts + gridlines
    try:
        if ctype not in (XL_CHART_TYPE.PIE,):
            for ax in (ch.category_axis, ch.value_axis):
                _font(ax.tick_labels.font, theme, theme.type_scale["caption"], theme.color("muted"))
            if not theme.chart_style.get("gridlines", False):
                ch.value_axis.has_major_gridlines = False
                ch.category_axis.has_major_gridlines = False
    except Exception:
        pass
    return gf


def _font(font, theme, size_pt, color):
    font.name = theme.font("body")
    font.size = Pt(size_pt)
    font.color.rgb = color
