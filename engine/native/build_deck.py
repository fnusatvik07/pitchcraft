#!/usr/bin/env python3
"""Build an editable .pptx from a brief.json + a theme.

Usage:
    python build_deck.py <brief.json> <theme> <out.pptx>

The brief.json `outline[]` drives the deck: each item's `layout` field selects a
builder from layouts.LAYOUTS; the rest of the item is its content. Output is a
fully native, editable PowerPoint (real text, shapes, tables, charts, pictures).
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# local imports (run from engine/ or with engine/ on path)
sys.path.insert(0, str(Path(__file__).resolve().parent))
import charts  # noqa: E402
from assets import place_image  # noqa: E402
from layouts import LAYOUTS  # noqa: E402
from theme import Theme  # noqa: E402

BLANK_LAYOUT = 6


def build(brief_path: str, theme_name: str, out_path: str) -> str:
    brief = json.loads(Path(brief_path).read_text())
    theme = Theme.load(theme_name)
    brief_dir = Path(brief_path).resolve().parent

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[BLANK_LAYOUT]

    ctx = {
        "place_image": lambda slide, img, x, y, w, h, th: place_image(
            slide, img, x, y, w, h, th, base_dir=brief_dir
        ),
        "add_chart": charts.add_chart,
    }

    outline = brief.get("outline", [])
    if not outline:
        raise ValueError("brief.json has no 'outline' items")

    for i, item in enumerate(outline):
        layout_key = item.get("layout", "content-image")
        builder = LAYOUTS.get(layout_key)
        if builder is None:
            raise ValueError(f"slide {i}: unknown layout '{layout_key}'. "
                             f"Known: {', '.join(sorted(LAYOUTS))}")
        slide = prs.slides.add_slide(blank)
        builder(slide, item, theme, ctx)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main(argv):
    if len(argv) != 4:
        print(__doc__)
        return 2
    out = build(argv[1], argv[2], argv[3])
    print(f"Built {out} ({_count(out)} slides) with theme '{argv[2]}'")
    return 0


def _count(path):
    return len(Presentation(path).slides._sldIdLst)


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
