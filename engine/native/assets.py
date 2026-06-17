"""Asset fetching + native embedding for PitchCraft.

- place_image(): embed a local image as a NATIVE picture shape, center-cropped to
  the target region. Never rasterizes the whole slide.
- ensure_icon(): fetch a brand logo / icon SVG from Iconify and rasterize to PNG
  via rsvg-convert (small content art only; the slide stays native).
- search_photo(): query the Openverse API for a CC0/CC-BY photo (used by the
  web-researcher agent), returning url + license + credit.

Licensing rule: only embed CC0 / CC-BY / public-domain media, or brand logos used
nominatively for identification. Every asset carries its license + credit.
"""

from __future__ import annotations

import subprocess
import urllib.parse
import urllib.request
from pathlib import Path

# NOTE: PIL and pptx are imported lazily inside the functions that need them.
# The asset-fetch path (ensure_icon / search_photo) is network-only and is used by
# the researcher agent; keeping pptx/lxml out of module import avoids a subprocess
# crash (lxml + fork) when calling rsvg-convert.

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
ICON_CACHE = PLUGIN_ROOT / "assets" / "icons"


def _resolve(path, base_dir: Path | None = None) -> Path | None:
    """Find an image given a path that may be absolute or relative to the cwd,
    the brief's dir, or the plugin root. Returns the first existing path."""
    p = Path(path)
    if p.is_absolute():
        return p if p.exists() else None
    candidates = [Path.cwd() / p]
    if base_dir:
        candidates.append(base_dir / p)
    candidates.append(PLUGIN_ROOT / p)
    for c in candidates:
        if c.exists():
            return c
    return None
ICONIFY = "https://api.iconify.design/{prefix}/{name}.svg"
OPENVERSE = "https://api.openverse.org/v1/images/"
UA = {"User-Agent": "PitchCraft/1.0 (+workshop deck generator)"}


# --------------------------------------------------------------------------- #
# embedding
# --------------------------------------------------------------------------- #
def place_image(slide, img, x, y, w, h, theme, base_dir: Path | None = None) -> bool:
    """img: dict with at least {path}. Returns True on success, False to fall back."""
    path = img.get("path") if isinstance(img, dict) else img
    if not path:
        return False
    p = _resolve(path, base_dir)
    if p is None:
        return False
    try:
        from PIL import Image
        from pptx.util import Inches
        im = Image.open(p)
        has_alpha = im.mode in ("RGBA", "LA", "P") and (
            im.mode != "P" or "transparency" in im.info
        )
        if has_alpha:
            # Logo / transparent art: CONTAIN-fit, centered, no crop, alpha preserved.
            iw, ih = im.size
            # leave room so a dark logo isn't edge-to-edge; reserve for a chip on dark themes
            box_w, box_h = w * 0.8, h * 0.8
            scale = min(box_w / iw, box_h / ih)
            nw, nh = iw * scale, ih * scale
            nx, ny = x + (w - nw) / 2, y + (h - nh) / 2
            # On dark themes a dark logo vanishes -> draw a light rounded chip behind it.
            bg = theme.palette.get("bg", "FFFFFF").lstrip("#")
            lum = 0.299 * int(bg[0:2], 16) + 0.587 * int(bg[2:4], 16) + 0.114 * int(bg[4:6], 16)
            if lum < 110:
                from layouts import add_rect  # lazy to avoid import cycle
                from theme import hex_to_rgb
                pad = 0.28
                add_rect(slide, nx - pad, ny - pad, nw + 2 * pad, nh + 2 * pad,
                         hex_to_rgb("FFFFFF"), rounded=True, radius_pt=14)
            slide.shapes.add_picture(str(p), Inches(nx), Inches(ny), Inches(nw), Inches(nh))
        else:
            # Photo: center-crop to fill the region exactly.
            cropped = _crop_to_aspect(p, w / h)
            slide.shapes.add_picture(str(cropped), Inches(x), Inches(y), Inches(w), Inches(h))
        if isinstance(img, dict) and img.get("credit"):
            from layouts import add_caption  # local import avoids cycle
            add_caption(slide, x, y + h + 0.03, w, theme, img["credit"])
        return True
    except Exception:
        return False


def _crop_to_aspect(path: Path, target_ratio: float) -> Path:
    """Center-crop a (non-transparent) image to target w/h ratio; write a sibling file."""
    from PIL import Image
    im = Image.open(path).convert("RGB")
    w, h = im.size
    cur = w / h
    if abs(cur - target_ratio) < 0.01:
        return path
    if cur > target_ratio:  # too wide -> crop width
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        im = im.crop((left, 0, left + new_w, h))
    else:  # too tall -> crop height
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        im = im.crop((0, top, w, top + new_h))
    out = path.with_name(path.stem + "_crop.png")
    im.save(out)
    return out


# --------------------------------------------------------------------------- #
# fetching (used by the web-researcher agent)
# --------------------------------------------------------------------------- #
def ensure_icon(name: str, prefix: str = "logos", height_px: int = 256) -> Path | None:
    """Fetch an Iconify icon/logo SVG and rasterize to PNG. Returns local path."""
    ICON_CACHE.mkdir(parents=True, exist_ok=True)
    png = ICON_CACHE / f"{prefix}_{name}.png"
    if png.exists():
        return png
    svg = ICON_CACHE / f"{prefix}_{name}.svg"
    url = ICONIFY.format(prefix=prefix, name=name)
    try:
        req = urllib.request.Request(url, headers=UA)
        data = urllib.request.urlopen(req, timeout=15).read()
        if b"<svg" not in data:
            return None
        svg.write_bytes(data)
        subprocess.run(
            ["rsvg-convert", "-h", str(height_px), "-o", str(png), str(svg)],
            check=True, capture_output=True,
        )
        return png if png.exists() else None
    except Exception:
        return None


def search_photo(query: str, license_str: str = "cc0,by", n: int = 3) -> list[dict]:
    """Query Openverse for license-filtered photos. Returns list of dicts:
    {url, thumbnail, license, credit, source}."""
    params = urllib.parse.urlencode({
        "q": query, "license": license_str, "page_size": n, "mature": "false",
    })
    try:
        req = urllib.request.Request(OPENVERSE + "?" + params, headers=UA)
        import json
        body = json.loads(urllib.request.urlopen(req, timeout=20).read())
    except Exception:
        return []
    out = []
    for r in body.get("results", []):
        creator = r.get("creator") or "Unknown"
        lic = (r.get("license") or "").upper()
        ver = r.get("license_version") or ""
        out.append({
            "url": r.get("url"),
            "thumbnail": r.get("thumbnail"),
            "license": f"CC {lic} {ver}".strip(),
            "credit": f"{r.get('title') or 'photo'} by {creator} (CC {lic} {ver})".strip(),
            "source": r.get("foreign_landing_url"),
        })
    return out


def download(url: str, dest: Path) -> Path | None:
    try:
        dest.parent.mkdir(parents=True, exist_ok=True)
        req = urllib.request.Request(url, headers=UA)
        dest.write_bytes(urllib.request.urlopen(req, timeout=30).read())
        return dest
    except Exception:
        return None


if __name__ == "__main__":
    import sys
    if len(sys.argv) >= 2 and sys.argv[1] == "icon":
        print(ensure_icon(sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else "logos"))
    elif len(sys.argv) >= 2 and sys.argv[1] == "photo":
        for r in search_photo(" ".join(sys.argv[2:])):
            print(r)
