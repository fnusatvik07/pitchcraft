#!/usr/bin/env python3
"""Render a folder of slide HTML files into a .pptx (full-bleed, retina).

The cinematic skill composes bespoke slide-NN.html files (linking cinematic.css);
this renders each at 1920x1080 @2x via Chromium and assembles them, in filename
order, into an editable-container .pptx (slides are images). Optional notes.json
( {"slide-01": "speaker note", ...} ) injects speaker notes.

Usage: python html_to_pptx.py <html_dir> <out.pptx> [notes.json]
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

W, H = 1920, 1080


def build(html_dir: str, out_path: str, notes_path: str | None = None):
    d = Path(html_dir)
    slides = sorted(d.glob("slide-*.html"))
    if not slides:
        slides = sorted(d.glob("*.html"))
    if not slides:
        raise SystemExit(f"no slide html files in {d}")
    notes = json.loads(Path(notes_path).read_text()) if notes_path and Path(notes_path).exists() else {}

    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Inches

    pngs = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        page = browser.new_page(viewport={"width": W, "height": H}, device_scale_factor=2)
        for hp in slides:
            page.goto(hp.resolve().as_uri(), wait_until="load")
            try:
                page.evaluate("async () => { await document.fonts.ready; }")
            except Exception:
                pass
            page.wait_for_timeout(250)  # let web fonts/images settle
            # auto-fit safety net: if any content overflows the 1920x1080 canvas,
            # uniformly shrink the .pad from center until it fits (never ship overflow)
            page.evaluate("""() => {
              const pad = document.querySelector('.pad'); if (!pad) return 1;
              const W=1920, H=1080, M=2;
              const overflow = () => {
                for (const el of pad.querySelectorAll('*')) {
                  const r = el.getBoundingClientRect();
                  if (r.width===0||r.height===0) continue;
                  if (r.bottom>H-M || r.right>W-M) return true;
                }
                return false;
              };
              if (!overflow()) return 1;
              pad.style.transformOrigin='50% 50%';
              let k=1;
              for (let i=0;i<22 && overflow(); i++){ k-=0.03; pad.style.transform=`scale(${k})`; }
              return k;
            }""")
            png = hp.with_suffix(".png")
            page.screenshot(path=str(png))
            pngs.append((hp.stem, png))
        browser.close()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for stem, png in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
        if stem in notes:
            s.notes_slide.notes_text_frame.text = notes[stem]
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Built {out_path} ({len(pngs)} slides) from {html_dir}/")
    return out_path


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__); raise SystemExit(2)
    build(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else None)
